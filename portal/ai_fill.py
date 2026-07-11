"""Gemini-backed auto-fill for analyst-call sections.

Admin uploads a pitch deck (.pptx); we extract its text and ask Gemini to draft the
headline summary (a short paragraph) plus a one-liner for each section. The API key
lives in st.secrets['gemini']['api_key'] (git-ignored) — it is NEVER hardcoded here,
because this is a public repo.

    [gemini]
    api_key = "..."
    # model = "gemini-2.0-flash"   # optional override
"""
import io
import json

import streamlit as st

_DEFAULT_MODEL = "gemini-2.0-flash"
_MAX_DECK_CHARS = 20000   # cap so a huge deck can't blow the request size


def _cfg() -> dict:
    try:
        return dict(st.secrets.get("gemini", {}) or {})
    except Exception:
        return {}


def gemini_ready() -> bool:
    """True when a Gemini API key is configured (enables the auto-fill button)."""
    return bool(_cfg().get("api_key"))


def extract_pptx_text(data: bytes, filename: str = "") -> str:
    """Pull the readable text (titles, bodies, tables, speaker notes) from a .pptx.

    Raises on legacy .ppt (python-pptx only reads the modern XML format) or an empty deck.
    """
    if filename.lower().endswith(".ppt") and not filename.lower().endswith(".pptx"):
        raise ValueError("Legacy .ppt isn't supported — please re-save the deck as .pptx.")
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("python-pptx isn't installed — add it to requirements.txt.") from e

    prs = Presentation(io.BytesIO(data))
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
            if getattr(shape, "has_table", False):
                try:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                except Exception:
                    pass
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note = slide.notes_slide.notes_text_frame.text.strip()
                if note:
                    parts.append(f"(notes) {note}")
        except Exception:
            pass
        if parts:
            chunks.append(f"--- Slide {i} ---\n" + "\n".join(parts))

    text = "\n\n".join(chunks).strip()
    if not text:
        raise ValueError("No readable text found in the deck.")
    return text


def fill_analyst_sections(deck_text: str, sections: list) -> dict:
    """Ask Gemini to draft {"summary": <paragraph>, <section>: <one-liner>, ...} as JSON.

    `summary` is a short paragraph; every section value is a single concise sentence.
    Returns only the known keys (summary + the given sections), all stripped strings.
    """
    import requests

    cfg = _cfg()
    api_key = cfg.get("api_key")
    if not api_key:
        raise RuntimeError("No Gemini API key configured (st.secrets['gemini']['api_key']).")
    model = cfg.get("model", _DEFAULT_MODEL)

    deck_text = (deck_text or "")[:_MAX_DECK_CHARS]
    section_lines = "\n".join(f'- "{s}": ONE concise sentence (<= 25 words).' for s in sections)
    prompt = (
        "You are a steel-market analyst assistant. Read the analyst-call pitch-deck text "
        "below and produce a briefing for a dashboard card.\n\n"
        "Return STRICT JSON only (no markdown fences, no commentary) with these keys:\n"
        '- "summary": a short headline paragraph of 2-4 sentences capturing the overall '
        "market view (prices, demand, direction).\n"
        f"{section_lines}\n\n"
        "Rules: base every line ONLY on the deck's content; do not invent numbers. If a "
        "section has no relevant content in the deck, return an empty string for it. Keep "
        "each section value to a single line of plain text (no bullet characters, no labels).\n\n"
        "=== DECK TEXT ===\n" + deck_text
    )

    schema = {
        "type": "object",
        "properties": {"summary": {"type": "string"},
                       **{s: {"type": "string"} for s in sections}},
        "required": ["summary"],
    }
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
    body = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {
            "temperature": 0.2,
            "responseMimeType": "application/json",
            "responseSchema": schema,
        },
    }
    # key in a header (not the URL query) so it never lands in request logs
    resp = requests.post(url, headers={"Content-Type": "application/json",
                                       "x-goog-api-key": api_key},
                         json=body, timeout=60)
    if resp.status_code != 200:
        raise RuntimeError(f"Gemini API error {resp.status_code}: {resp.text[:300]}")

    data = resp.json()
    try:
        text = data["candidates"][0]["content"]["parts"][0]["text"]
    except (KeyError, IndexError) as e:
        raise RuntimeError(f"Unexpected Gemini response: {json.dumps(data)[:300]}") from e

    try:
        out = json.loads(text)
    except json.JSONDecodeError as e:
        raise RuntimeError(f"Gemini did not return valid JSON: {text[:300]}") from e

    result = {"summary": str(out.get("summary", "")).strip()}
    for s in sections:
        result[s] = str(out.get(s, "")).strip()
    return result
